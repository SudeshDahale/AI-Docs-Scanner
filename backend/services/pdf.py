import io
from pypdf import PdfReader

def extract_text(file, filename: str = ""):
    """
    Returns list of {page: int, text: str} for any supported file type.
    Supported: PDF, DOCX, DOC, PPTX, PPT, TXT, MD
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    file_bytes = file.read()

    if ext == "pdf":
        return _extract_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return _extract_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return _extract_pptx(file_bytes)
    elif ext in ("txt", "md"):
        return _extract_txt(file_bytes)
    else:
        # Fallback: try PDF first, then plain text
        try:
            return _extract_pdf(file_bytes)
        except Exception:
            return _extract_txt(file_bytes)


def _extract_pdf(file_bytes):
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"page": i + 1, "text": text})
    return pages or [{"page": 1, "text": "No text could be extracted from this PDF."}]


def _extract_docx(file_bytes):
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    pages = []
    current_page = 1
    current_text = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Treat every ~40 paragraphs as a "page" (Word has no real page concept)
        current_text.append(text)
        if len(current_text) >= 40:
            pages.append({"page": current_page, "text": "\n".join(current_text)})
            current_page += 1
            current_text = []

    if current_text:
        pages.append({"page": current_page, "text": "\n".join(current_text)})

    return pages or [{"page": 1, "text": "No text could be extracted from this Word document."}]


def _extract_pptx(file_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({"page": i + 1, "text": "\n".join(texts)})

    return pages or [{"page": 1, "text": "No text could be extracted from this presentation."}]


def _extract_txt(file_bytes):
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1", errors="replace")

    # Split into ~2000 char pages
    pages = []
    chunk_size = 2000
    for i in range(0, len(text), chunk_size):
        chunk = text[i:i + chunk_size].strip()
        if chunk:
            pages.append({"page": i // chunk_size + 1, "text": chunk})

    return pages or [{"page": 1, "text": "File appears to be empty."}]